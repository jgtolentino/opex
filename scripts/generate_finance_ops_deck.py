#!/usr/bin/env python3
"""
Generate Finance Operations Overview PowerPoint deck.

Requirements:
    pip install python-pptx

Usage:
    python scripts/generate_finance_ops_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Color scheme
COLOR_PRIMARY = RGBColor(0, 51, 102)  # Navy blue
COLOR_ACCENT = RGBColor(0, 176, 80)  # Green
COLOR_WARNING = RGBColor(255, 192, 0)  # Orange
COLOR_TEXT = RGBColor(64, 64, 64)  # Dark gray

def create_title_slide(prs, title, subtitle):
    """Create title slide."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY

    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT

    return slide

def create_content_slide(prs, title, bullet_points):
    """Create content slide with bullet points."""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]

    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY

    text_frame = content_shape.text_frame
    text_frame.clear()

    for bullet in bullet_points:
        p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT

    return slide

def create_two_column_slide(prs, title, left_content, right_content):
    """Create slide with two columns."""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(4.5), Inches(5)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT

    # Right column
    right_box = slide.shapes.add_textbox(
        Inches(5), Inches(1.5), Inches(4.5), Inches(5)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT

    return slide

def main():
    """Generate Finance Operations Overview deck."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    create_title_slide(
        prs,
        "Finance Operations Overview",
        "Standardized Month-End Close & Tax Compliance Framework\nFinance Shared Services Center"
    )

    # Slide 2: Current Challenges
    create_content_slide(
        prs,
        "Current Pain Points",
        [
            "❌ Manual, inconsistent close processes across 8 agencies",
            "⚠️ Late month-end close (15-20 days average)",
            "📊 High error rates due to manual data entry and reconciliation",
            "🚨 BIR filing deadline pressure and compliance risk",
            "👥 Finance team overload during close period (94% high workload)",
            "📝 Poor documentation and knowledge loss during turnover",
            "🔄 Repetitive manual tasks consuming 60%+ of staff time"
        ]
    )

    # Slide 3: New Standardized Framework
    create_content_slide(
        prs,
        "New Standardized Close & Tax Framework",
        [
            "✅ 15-day structured close cycle (Day 1-15 following month-end)",
            "📋 47 standardized tasks with clear ownership (RACI)",
            "🗓️ Integrated BIR tax filing calendar (monthly & quarterly)",
            "🤖 85% automation target for routine reconciliations",
            "📊 Real-time dashboards and progress tracking",
            "🛡️ Built-in validation gates and quality controls",
            "📚 Comprehensive policies and SOPs documented"
        ]
    )

    # Slide 4: Month-End Calendar
    create_two_column_slide(
        prs,
        "Month-End Close Calendar",
        [
            "📅 Day 1-2: Bank Reconciliation",
            "  • Download statements",
            "  • Match transactions",
            "  • Investigate discrepancies",
            "",
            "📅 Day 3-4: AP/AR Reconciliation",
            "  • Generate aging reports",
            "  • Reconcile sub-ledgers to GL",
            "  • Review aged items",
            "",
            "📅 Day 5-6: Fixed Assets & Depreciation",
            "  • Update asset register",
            "  • Calculate depreciation",
            "  • Post JEs",
            "",
            "📅 Day 7-8: Inventory & COGS",
            "  • Physical inventory count",
            "  • Reconcile to system",
            "  • Post adjustments"
        ],
        [
            "📅 Day 9-10: Journal Entries & Accruals",
            "  • Recurring JEs",
            "  • Accrual entries",
            "  • Prepayment amortization",
            "  • 🏛️ BIR Form 1601-C filing",
            "",
            "📅 Day 11-12: Financial Statements",
            "  • Trial balance",
            "  • Income statement",
            "  • Balance sheet",
            "  • Cash flow statement",
            "",
            "📅 Day 13-14: Review & Approval",
            "  • SFM review",
            "  • Address findings",
            "  • FD final approval",
            "",
            "📅 Day 15: Period Lock & Archive",
            "  • Lock period in ERP",
            "  • Archive documentation"
        ]
    )

    # Slide 5: RACI Overview
    create_two_column_slide(
        prs,
        "RACI - Roles & Responsibilities",
        [
            "🎯 Finance SSC Manager (FD)",
            "  • Strategic oversight",
            "  • Final financial statement approval",
            "  • Policy compliance",
            "  • Automation sponsor",
            "",
            "👔 Senior Finance Manager (SFM)",
            "  • Review & approve reconciliations",
            "  • Review & approve journal entries",
            "  • BIR filing oversight",
            "  • Financial statement review",
            "  • Period lock authority"
        ],
        [
            "💼 Finance Supervisor (FS)",
            "  • Day-to-day execution (27 tasks)",
            "  • Prepare JEs and reconciliations",
            "  • BIR form preparation & filing",
            "  • Financial statement preparation",
            "",
            "👤 Accounting Staff (AS)",
            "  • Data entry and processing (10 tasks)",
            "  • Report generation from ERP",
            "  • Document archiving",
            "  • Cash advance processing",
            "  • Reconciliation support"
        ]
    )

    # Slide 6: BIR Compliance Map
    create_content_slide(
        prs,
        "BIR Tax Filing Compliance Map",
        [
            "📋 Monthly Forms (Day 10 deadline):",
            "  • Form 1601-C: Compensation withholding tax",
            "  • Form 0619-E/F: Final income tax withheld",
            "",
            "📋 Quarterly Forms (Day 25 deadline):",
            "  • Form 2550Q: Quarterly VAT return",
            "  • Form 1701-Q: Expanded withholding tax",
            "  • Form 1702-Q: Quarterly income tax return",
            "",
            "🔄 4-Step Filing Workflow:",
            "  1. Preparation (BIR Deadline – 4 days): Finance Supervisor",
            "  2. Report Approval (BIR Deadline – 2 days): Senior Finance Manager",
            "  3. Payment Approval (BIR Deadline – 1 day): Finance Director",
            "  4. Filing & Payment (On or before deadline): Finance Supervisor via eFPS"
        ]
    )

    # Slide 7: Automation Opportunities
    create_two_column_slide(
        prs,
        "Automation Opportunities & Impact",
        [
            "🤖 Bank Reconciliation (Day 1-2)",
            "  • Automatic transaction import",
            "  • Rule-based matching (85%+ auto-match)",
            "  • Variance alerts",
            "  Impact: 6h → 1.5h (75% reduction)",
            "",
            "🤖 Depreciation Calculation (Day 5-6)",
            "  • Auto-calculate from asset register",
            "  • Generate JEs automatically",
            "  • Schedule validation",
            "  Impact: 4h → 0.5h (87% reduction)",
            "",
            "🤖 Accrual Reversals (Day 9)",
            "  • Automatic reversal on Day 1",
            "  • Prepayment amortization",
            "  • Recurring entry automation",
            "  Impact: 2h → 0.25h (87% reduction)"
        ],
        [
            "🤖 VAT Reconciliation (Day 9-10)",
            "  • Input/output VAT matching",
            "  • Summary report generation",
            "  • Variance detection",
            "  Impact: 5h → 1h (80% reduction)",
            "",
            "🤖 BIR eFPS Integration (Day 10)",
            "  • Auto-populate from ERP data",
            "  • Electronic filing submission",
            "  • Payment confirmation tracking",
            "  Impact: 3h → 0.5h (83% reduction)",
            "",
            "📊 Overall Impact:",
            "  • 60% of manual tasks automated",
            "  • Close cycle: 20 days → 12 days",
            "  • Error rate: 5% → <1%",
            "  • Staff capacity freed: 40%"
        ]
    )

    # Slide 8: Implementation Phases
    create_content_slide(
        prs,
        "Implementation Roadmap",
        [
            "🔷 Phase 1: Foundation (Month 1-2) - Documentation & Training",
            "  • Finalize policies and SOPs",
            "  • Train all finance staff on new process",
            "  • Set up Notion task tracking databases",
            "",
            "🔷 Phase 2: Quick Wins (Month 2-3) - Priority Automation",
            "  • Bank reconciliation automation",
            "  • Depreciation calculation automation",
            "  • Accrual reversal automation",
            "",
            "🔷 Phase 3: BIR Integration (Month 3-4) - Compliance Automation",
            "  • VAT reconciliation automation",
            "  • BIR eFPS integration stub",
            "  • Tax filing workflow automation",
            "",
            "🔷 Phase 4: Optimization (Month 4-6) - Continuous Improvement",
            "  • Dashboard and reporting automation",
            "  • Process refinement based on metrics",
            "  • Advanced analytics and forecasting"
        ]
    )

    # Slide 9: Success Metrics
    create_two_column_slide(
        prs,
        "Success Criteria & KPIs",
        [
            "📊 Primary Metrics:",
            "  • Close completion: ≤12 days (vs 20 current)",
            "  • BIR filing timeliness: 100% on-time",
            "  • Reconciliation accuracy: 100%",
            "  • Automation rate: ≥85%",
            "  • Error rate: <1% (vs 5% current)",
            "",
            "📊 Efficiency Metrics:",
            "  • Staff time saved: 40%",
            "  • Manual task reduction: 60%",
            "  • Review turnaround: <4 hours",
            "  • Documentation completeness: 100%"
        ],
        [
            "📊 Quality Metrics:",
            "  • Zero BIR penalties",
            "  • Zero period reopening",
            "  • Zero audit findings",
            "  • 100% policy compliance",
            "",
            "📊 Team Metrics:",
            "  • Team utilization: 80-90%",
            "  • Overtime reduction: 50%",
            "  • Training completion: 100%",
            "  • Employee satisfaction: +20%"
        ]
    )

    # Slide 10: Risk Mitigation
    create_content_slide(
        prs,
        "Risk Management & Mitigation",
        [
            "⚠️ Change Management Risk",
            "  → Mitigation: Phased rollout, comprehensive training, champions in each team",
            "",
            "⚠️ System/Automation Failure Risk",
            "  → Mitigation: Manual override procedures, error monitoring, daily health checks",
            "",
            "⚠️ Regulatory Change Risk",
            "  → Mitigation: Quarterly policy review, BIR update monitoring, flexible workflows",
            "",
            "⚠️ Resource Constraint Risk",
            "  → Mitigation: Cross-training, documentation, automation priorities, phased approach",
            "",
            "⚠️ Data Quality Risk",
            "  → Mitigation: Validation gates, error alerts, reconciliation controls, audit trails"
        ]
    )

    # Slide 11: Technology Stack
    create_two_column_slide(
        prs,
        "Technology & Systems",
        [
            "💻 Core Systems:",
            "  • ERP: Primary financial system",
            "  • Supabase PostgreSQL: Automation backend",
            "  • n8n: Workflow automation engine",
            "  • Notion: Task tracking and knowledge base",
            "",
            "🔗 Integrations:",
            "  • ERP ↔ Supabase: Real-time data sync",
            "  • Supabase ↔ n8n: Automation triggers",
            "  • n8n ↔ BIR eFPS: Electronic filing",
            "  • Notion ↔ Supabase: Task status sync"
        ],
        [
            "🛠️ Automation Workflows:",
            "  • Bank statement import (n8n)",
            "  • Transaction matching (Supabase functions)",
            "  • Depreciation calculation (SQL procedures)",
            "  • Journal entry generation (n8n workflows)",
            "  • BIR form population (Supabase → eFPS)",
            "",
            "📊 Reporting & Dashboards:",
            "  • Real-time close progress dashboard",
            "  • KPI tracking (Notion/Supabase)",
            "  • BIR filing status tracker",
            "  • Exception and error monitoring"
        ]
    )

    # Slide 12: Investment & ROI
    create_two_column_slide(
        prs,
        "Investment & Return on Investment",
        [
            "💰 Initial Investment:",
            "  • Documentation: 40 hours (complete)",
            "  • Training: 80 hours (2 weeks)",
            "  • Automation development: 160 hours (1 month)",
            "  • Testing & validation: 40 hours",
            "  Total: 320 hours (~2 months FTE)",
            "",
            "💸 Ongoing Costs:",
            "  • Notion subscription: ~$10/month",
            "  • Supabase Pro: $25/month",
            "  • n8n hosting: $20/month",
            "  • Maintenance: 10 hours/month",
            "  Total: ~$55/month + 10h maintenance"
        ],
        [
            "📈 Monthly Benefits:",
            "  • Time saved: 80 hours/month (40% capacity)",
            "  • Error reduction: ~20 hours rework avoided",
            "  • Overtime elimination: 40 hours saved",
            "  • BIR penalty avoidance: Priceless",
            "  Total: 140 hours/month saved",
            "",
            "🎯 ROI Calculation:",
            "  • Break-even: 2.3 months",
            "  • Annual savings: 1,680 hours = 1 FTE",
            "  • Cost avoidance: Zero BIR penalties",
            "  • Strategic value: Data-driven insights",
            "  Payback period: <3 months"
        ]
    )

    # Slide 13: Next Steps
    create_content_slide(
        prs,
        "Immediate Next Steps",
        [
            "✅ Week 1-2: Foundation Setup",
            "  → Import Notion databases (task list, BIR schedule, policies)",
            "  → Conduct team orientation and training sessions",
            "  → Set up Supabase backend infrastructure",
            "",
            "✅ Week 3-4: Quick Win Automation",
            "  → Deploy bank reconciliation automation",
            "  → Deploy depreciation calculation automation",
            "  → Validate with dry-run month-end close",
            "",
            "✅ Week 5-6: BIR Integration",
            "  → Implement BIR eFPS integration stub",
            "  → Test with January 2026 Form 1601-C filing",
            "  → Deploy VAT reconciliation automation",
            "",
            "✅ Week 7-8: Full Rollout",
            "  → Execute first full month-end close with new framework",
            "  → Measure KPIs and document lessons learned",
            "  → Continuous improvement based on metrics"
        ]
    )

    # Slide 14: Questions & Discussion
    create_title_slide(
        prs,
        "Questions & Discussion",
        "Finance SSC Manager\nFinance Shared Services Center"
    )

    # Save presentation
    output_path = "exec/Finance_Operations_Overview.pptx"
    prs.save(output_path)
    print(f"✅ PowerPoint deck created: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print("\n📋 Slide breakdown:")
    for i, slide in enumerate(prs.slides, 1):
        title = slide.shapes.title.text if slide.shapes.title else "(Title slide)"
        print(f"  {i}. {title}")

if __name__ == "__main__":
    main()
